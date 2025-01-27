import json
import os
from concurrent.futures import ThreadPoolExecutor
from glob import glob
from typing import Literal

import func_argparse
import jsonlines
from jinja2 import Template
from PIL import Image
from pptx import Presentation
from torch import cosine_similarity
from tqdm import tqdm
from transformers import CLIPModel, CLIPProcessor

import llms
from presentation import Presentation
from utils import edit_distance, ppt_to_images

outline_template = Template(
    """
    From the following text which contains a set of headings and some content within each heading:
    {{ text }}
    Extract the most important headings present in it. Reduce the length of each heading to five words if they are lengthy.
    Example Output:
    ["Heading 1", "Heading 2", "Heading 3"]
    Output: give your output as a list of strings in json format
    """
)
mapping_template = Template(
    """
    Think step by step and then answer the following question:  You are given with the following title: {{outline_headings}}
    and a list of keys: {{document_heading_from_bird_eye_view}}
    Each key is associated with some text as presented in the dictionary format below:
    {{bird_eye_view}}
    The task is to find 1-2 significantly matched keys. The matching should be done based on the similarity of the text associated with the keys with the given heading.
    Example Output:
    thoughts...
    {"Heading 1": ["key1", "key2"], "Heading 2": ["key1", "key4"]}
    Output: give your final output as a dictionary in json format, notice that all headings must be present in the output, no heading should be left out and at least one key should be present in the output for each heading
    """
)
generation_template = Template(
    """
    You are a presentation generator from a source of text. You have to generate the slide number {{slide_index}}. Previous slide headings and slide contents are given below in the format of a list of dictionaries. {{previous_slide}} Given the following slide heading and the source of text respectively, create the content of the slide number {{slide_index}} such that: 1. The slide should have maximum {{max_bullet}} bullet points. 2. Ensure that the content of the bullet points are coming strictly from the given source of text only. 3. The content of the slide is very relevant to the given slide heading 4. Each bullet point should have a maximum of 10 words 5. Ensure that this slide does not have any content repeated from the previous slides. 6. The flow of the overall presentation is nice. 7. Do not prefix the slide title before the bullet poide nts in the output  SliTitle: {{slide_heading}} Source of text: {{text}}
    Example Output:
    ["bullet point 1", "bullet point 2"]
    Output: give your output as a list of strings in json format
    """
)


def filter_aspect_ratio(image: list[str]):
    filtered_images = []
    for i in image:
        size = image = Image.open(i).size
        long, short = max(size), min(size)
        if long / short < 4:
            filtered_images.append(i)
    return filtered_images


def get_indexed_sections(bird_eye: dict, indexs: list[str]):
    indexed_sections = []
    for section in bird_eye["sections"]:
        for subsection in section["subsections"]:
            if any(edit_distance(key, next(iter(subsection))) > 0.9 for key in indexs):
                indexed_sections.append(subsection)
    return indexed_sections


def generate_content(source_text: str, bird_eye: dict, max_bullet: int):
    bird_eye_headdings = []
    for section in bird_eye["sections"]:
        bird_eye_headdings.extend(
            [next(iter(subsec)) for subsec in section["subsections"]]
        )
    outline: list[str] = llms.language_model(
        outline_template.render(text=source_text), return_json=True
    )
    assert len(outline) != 0, "No outline found"
    mapping = llms.language_model(
        mapping_template.render(
            outline_headings=outline,
            document_heading_from_bird_eye_view=bird_eye_headdings,
            bird_eye_view=bird_eye,
        ),
        return_json=True,
    )
    slides = []
    for slide_title in outline:
        bullet_points = llms.language_model(
            generation_template.render(
                slide_heading=slide_title,
                text=get_indexed_sections(bird_eye, mapping.get(slide_title, [])),
                previous_slide=slides,
                max_bullet=max_bullet,
            ),
            return_json=True,
        )
        slides.append(
            {
                "title": slide_title,
                "bullets": bullet_points,
                "indexed_sections": mapping.get(slide_title, []),
            }
        )
    return slides


def generate_slides(
    output_dir: str,
    source_text: str,
    bird_eye: dict,
    images: list[str],
    model: CLIPModel,
    processor: CLIPProcessor,
):
    os.makedirs(output_dir, exist_ok=True)
    images = filter_aspect_ratio(images)
    slides = generate_content(source_text, bird_eye, 7)
    image_embeddings = model.get_image_features(
        **processor(images=[Image.open(i) for i in images], return_tensors="pt").to(
            "cuda"
        )
    ).unsqueeze(0)
    text_embeddings = model.get_text_features(
        **processor(
            text=["\n".join(slide["bullets"]) for slide in slides],
            return_tensors="pt",
            padding=True,
            max_length=77,
            truncation=True,
        ).to("cuda")
    ).unsqueeze(1)
    similarity = cosine_similarity(image_embeddings, text_embeddings, dim=-1)
    pptx = Presentation()
    for slide_idx, slide in enumerate(slides):  # match image here
        title = slide["title"]
        bullets = slide["bullets"]

        subsimilarity = similarity[slide_idx]
        if subsimilarity.max() > 0.8:
            slide = pptx.slides.add_slide(pptx.slide_layouts[6])
            bullets_placeholder = slide.shapes.placeholders[2]
            image = images[subsimilarity.argmax()]
            slides[slide_idx]["image"] = image
            slide.shapes.placeholders[1].insert_picture(image)
        else:
            slide = pptx.slides.add_slide(pptx.slide_layouts[1])
            bullets_placeholder = slide.shapes.placeholders[1]
        slide.shapes.title.text = title
        text_frame = bullets_placeholder.text_frame
        for bullet in bullets:
            para = text_frame.add_paragraph()
            para.text = bullet
            para.level = 1
    with jsonlines.open(output_dir + "/final.jsonl", "w") as writer:
        writer.write_all(slides)
    pptx.save(output_dir + "/final.pptx")
    ppt_to_images(output_dir + "/final.pptx", output_dir + "/slide_images")


def generate(model: Literal["Qwen2.5", "gpt"]):
    if model == "Qwen2.5":
        llms.language_model = llms.qwen2_5
    elif model == "gpt":
        llms.language_model = llms.gpt4o

    print("Generating slides on baseline with ", llms.language_model.model)
    llm_name = llms.get_simple_modelname(llms.language_model)
    model = CLIPModel.from_pretrained("openai/clip-vit-large-patch14").to("cuda").eval()
    processor = CLIPProcessor.from_pretrained("openai/clip-vit-large-patch14")
    folders = list(glob("data/*/pdf/*"))
    progress = tqdm(total=len(folders))

    def process_folder(pdf_folder, model, processor):
        source_text = open(f"{pdf_folder}/source.md").read()
        bird_eye = json.load(open(f"{pdf_folder}/refined_doc.json"))
        images = json.load(open(f"{pdf_folder}/image_caption.json")).keys()
        output_dir = f"{pdf_folder}/docpres/{llm_name}"
        if os.path.exists(output_dir + "/final.jsonl"):
            progress.write(f"Skipping {pdf_folder}")
            progress.update(1)
            return
        try:
            generate_slides(
                output_dir,
                source_text,
                bird_eye,
                list(images),
                model,
                processor,
            )
            progress.update(1)
        except Exception as e:
            print(f"Error in {pdf_folder}: {e}")

    # for folder in folders:
    #     process_folder(folder, model, processor)

    with ThreadPoolExecutor() as executor:
        list(executor.map(lambda f: process_folder(f, model, processor), folders))


if __name__ == "__main__":
    func_argparse.main([generate])
