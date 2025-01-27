import re
import traceback
from dataclasses import dataclass
from typing import Callable

from pptx import Presentation as PPTXPre
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape as PPTXAutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector as PPTXConnector
from pptx.shapes.group import GroupShape as PPTXGroupShape
from pptx.shapes.picture import Picture as PPTXPicture
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from pptx.slide import Slide as PPTXSlide
from pptx.text.text import _Paragraph, _Run
from rich import print

from utils import (
    IMAGE_EXTENSIONS,
    Config,
    apply_fill,
    dict_to_object,
    extract_fill,
    get_font_pptcstyle,
    get_font_style,
    merge_dict,
    object_to_dict,
    parse_groupshape,
    pexists,
    pjoin,
    runs_merge,
    wmf_to_images,
)

INDENT = "\t"


# textframe: shape bounds font
# paragraph: space, alignment, level, font bullet
# run: font, hyperlink, text
@dataclass
class StyleArg:
    """
    A class to represent style arguments for HTML conversion.
    """

    paragraph_id: bool = True
    element_id: bool = True
    font_style: bool = True
    area: bool = False
    size: bool = False
    geometry: bool = False
    show_image: bool = True


@dataclass
class Closure:
    closure: Callable
    paragraph_id: int = -1

    def apply(self, shape: BaseShape):
        """
        Apply the closure to a shape.

        Args:
            shape (BaseShape): The shape to apply the closure to.
        """
        self.closure(shape)

    def __gt__(self, other):
        if self.paragraph_id != other.paragraph_id:
            return self.paragraph_id > other.paragraph_id


class Paragraph:
    def __init__(self, paragraph: _Paragraph, idx: int):
        run = runs_merge(paragraph)
        self.idx = idx
        self.real_idx = idx
        self.bullet = paragraph.bullet
        if run is None:
            self.idx = -1
            return
        self.font = merge_dict(
            object_to_dict(paragraph.font), [object_to_dict(run.font)]
        )
        self.text = re.sub(r"(_x000B_|\\x0b)", " ", paragraph.text)

    def to_html(self, style_args: StyleArg):
        if self.idx == -1:
            raise ValueError(f"paragraph {self.idx} is not valid")
        tag = "li" if self.bullet else "p"
        id_str = f" id='{self.idx}'" if style_args.paragraph_id else ""
        font_style = get_font_style(self.font)
        style_str = (
            f" style='{font_style}'" if style_args.font_style and font_style else ""
        )
        if self.bullet:
            style_str += f" bullet-type='{self.bullet}'"
        return f"<{tag}{id_str}{style_str}>{self.text}</{tag}>"

    def __repr__(self):
        return f"Paragraph-{self.idx}: {self.text}"


class TextFrame:
    def __init__(self, shape: BaseShape, level: int):
        if not shape.has_text_frame:
            self.is_textframe = False
            return
        self.paragraphs = [
            Paragraph(paragraph, idx)
            for idx, paragraph in enumerate(shape.text_frame.paragraphs)
        ]
        para_offset = 0
        for para in self.paragraphs:
            if para.idx == -1:
                para_offset += 1
            else:
                para.idx = para.idx - para_offset
        if len(self.paragraphs) == 0:
            self.is_textframe = False
            return
        self.level = level
        self.text = shape.text
        self.is_textframe = True
        self.font = merge_dict(
            object_to_dict(shape.text_frame.font),
            [para.font for para in self.paragraphs if para.idx != -1],
        )

    def to_html(self, style_args: StyleArg):
        """
        Convert the text frame to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the text frame.
        """
        if not self.is_textframe:
            return ""
        repr_list = [
            para.to_html(style_args) for para in self.paragraphs if para.idx != -1
        ]
        return "\n".join([INDENT * self.level + repr for repr in repr_list])

    def __repr__(self):
        if not self.is_textframe:
            return "TextFrame: null"
        return f"TextFrame: {self.paragraphs}"

    def __len__(self):
        if not self.is_textframe:
            return 0
        return len(self.text)

    def to_pptc(self, father_idx: int) -> str:
        """
        Convert the text frame to PPTC format.

        Args:
            father_idx (int): The index of the parent shape.

        Returns:
            str: The PPTC representation of the text frame.
        """
        if not self.is_textframe:
            return ""
        s = f"[Text id={father_idx}]"
        for para in self.paragraphs:
            if para.idx == -1:
                continue
            s += f"\n"
            s += f"[Paragraph id={para.idx}]"
            s += get_font_pptcstyle(para.font) + f"\n"
            s += para.text + "\n"
        return s


class ShapeElement:
    def __init__(
        self,
        slide_idx: int,
        shape_idx: int,
        style: dict,
        data: dict,
        text_frame: TextFrame,
        slide_area: float,
        level: int,
    ):
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.style = style
        self.data = data
        self.text_frame = text_frame
        self._closure_keys = ["clone", "replace", "delete", "style"]
        self._closures: dict[str, list[Closure]] = {
            key: [] for key in self._closure_keys
        }
        self.slide_area = slide_area
        self.level = level

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        config: Config,
        slide_area: float,
        level: int = 0,
    ):
        """
        Create a ShapeElement from a BaseShape.

        Args:
            slide_idx (int): The index of the slide.
            shape_idx (int): The index of the shape.
            shape (BaseShape): The shape object.
            config (Config): The configuration object.
            slide_area (float): The area of the slide.
            level (int): The indentation level.

        Returns:
            ShapeElement: The created ShapeElement.
        """
        if shape_idx > 100 and isinstance(shape, PPTXGroupShape):
            raise ValueError(f"nested group shapes are not allowed")
        line = None
        if "line" in dir(shape) and shape.line._ln is not None:
            line = {
                "fill": extract_fill(shape.line),
                "width": shape.line.width,
                "dash_style": shape.line.dash_style,
            }
        fill = extract_fill(shape)
        style = {
            "shape_bounds": {
                "width": shape.width,
                "height": shape.height,
                "left": shape.left,
                "top": shape.top,
            },
            "shape_type": str(shape.shape_type).split("(")[0].lower(),
            "rotation": shape.rotation,
            "fill": fill,
            "line": line,
        }
        text_frame = TextFrame(shape, level + 1)
        obj = SHAPECAST.get(shape.shape_type, UnsupportedShape).from_shape(
            slide_idx,
            shape_idx,
            shape,
            style,
            text_frame,
            config,
            slide_area,
            level,
        )
        obj.xml = shape._element.xml
        # ? for debug, mask to enable pickling
        # obj.shape = shape
        return obj

    def build(self, slide: PPTXSlide):
        """
        Build the shape element in a slide.

        Args:
            slide (PPTXSlide): The slide to build the shape in.

        Returns:
            The built shape.
        """
        return slide.shapes._shape_factory(
            slide.shapes._spTree.insert_element_before(parse_xml(self.xml), "p:extLst")
        )

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}: shape {self.shape_idx} of slide {self.slide_idx}"

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the shape element to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the shape element.
        """
        return ""

    @property
    def closures(self):
        """
        Get the closures associated with the shape element.

        Returns:
            list: A list of closures.
        """
        closures = []
        closures.extend(sorted(self._closures["clone"]))
        closures.extend(self._closures["replace"] + self._closures["style"])
        closures.extend(sorted(self._closures["delete"], reverse=True))
        return closures

    @property
    def indent(self):
        return "\t" * self.level

    @property
    def left(self):
        return self.style["shape_bounds"]["left"].pt

    @left.setter
    def left(self, value):
        self.style["shape_bounds"]["left"] = value

    @property
    def top(self):
        return self.style["shape_bounds"]["top"].pt

    @top.setter
    def top(self, value):
        self.style["shape_bounds"]["top"] = value

    @property
    def width(self):
        return self.style["shape_bounds"]["width"].pt

    @width.setter
    def width(self, value):
        self.style["shape_bounds"]["width"] = value

    @property
    def height(self):
        return self.style["shape_bounds"]["height"].pt

    @height.setter
    def height(self, value):
        self.style["shape_bounds"]["height"] = value

    @property
    def area(self):
        """
        Get the area of the shape element.

        Returns:
            float: The area in square points.
        """
        return self.width * self.height

    @property
    def pptc_text_info(self):
        """
        Get the PPTC text information of the shape element.

        Returns:
            str: The PPTC text information.
        """
        if isinstance(self, Picture):
            return self.caption
        return self.text_frame.to_pptc(self.shape_idx)

    @property
    def pptc_space_info(self):
        """
        Get the PPTC space information of the shape element.

        Returns:
            str: The PPTC space information.
        """
        return f"Visual Positions: left={self.left}pt, top={self.top}pt\n"

    @property
    def pptc_size_info(self):
        """
        Get the PPTC size information of the shape element.

        Returns:
            str: The PPTC size information.
        """
        return f"Size: height={self.height}pt, width={self.width}pt\n"

    @property
    def pptc_description(self):
        """
        Get the PPTC description of the shape element.

        Returns:
            str: The PPTC description.
        """
        return f"[{self.__class__.__name__} id={self.shape_idx}]\n"

    def to_pptc(self):
        """
        Convert the shape element to PPTC format.

        Returns:
            str: The PPTC representation of the shape element.
        """
        s = ""
        s += self.pptc_description
        s += self.pptc_size_info
        s += self.pptc_space_info
        s += self.pptc_text_info
        return s

    def get_inline_style(self, style_args: StyleArg):
        """
        Get the inline style for the shape element.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The inline style string.
        """
        id_str = f" id='{self.shape_idx}'" if style_args.element_id else ""
        styles = []
        if style_args.area:
            styles.append(f"data-relative-area={self.area*100/self.slide_area:.2f}%;")
        if style_args.size:
            styles.append(f"width: {self.width}pt; height: {self.height}pt;")
        if style_args.geometry:
            styles.append(f"left: {self.left}pt; top: {self.top}pt;")
        if style_args.font_style and self.text_frame.is_textframe:
            font_style = get_font_style(self.text_frame.font)
            if font_style:
                styles.append(font_style)
        if len(styles) != 0:
            return id_str + " style='" + " ".join(styles) + "'"
        return id_str


class UnsupportedShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        *args,
        **kwargs,
    ):
        raise ValueError(f"unsupported shape {shape.shape_type}")


class TextBox(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: TextFrame,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        return cls(slide_idx, shape_idx, style, None, text_frame, slide_area, level)

    def to_html(self, style_args: StyleArg) -> str:

        return (
            f"{self.indent}<div{self.get_inline_style(style_args)}>\n"
            + self.text_frame.to_html(style_args)
            + f"\n{self.indent}</div>\n"
        )


class Picture(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXPicture,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        img_path = pjoin(
            config.IMAGE_DIR,
            f"{shape.image.sha1}.{shape.image.ext}",
        )
        if shape.image.ext == "wmf":
            img_path = img_path.replace(".wmf", ".jpg")
            if not pexists(img_path):
                wmf_to_images(shape.image.blob, img_path)
        elif shape.image.ext not in IMAGE_EXTENSIONS:
            raise ValueError(f"unsupported image type {shape.image.ext}")
        if not pexists(img_path):
            with open(img_path, "wb") as f:
                f.write(shape.image.blob)
        style["img_style"] = {
            "crop_bottom": shape.crop_bottom,
            "crop_top": shape.crop_top,
            "crop_left": shape.crop_left,
            "crop_right": shape.crop_right,
        }
        picture = cls(
            slide_idx,
            shape_idx,
            style,
            [img_path, shape.name, ""],
            text_frame,
            slide_area,
            level=level,
        )
        return picture

    def build(self, slide: PPTXSlide):
        shape = slide.shapes.add_picture(
            self.img_path,
            **self.style["shape_bounds"],
        )
        shape.name = self.data[1]
        dict_to_object(self.style["img_style"], shape.image)
        apply_fill(shape, self.style["fill"])
        if self.style["line"] is not None:
            apply_fill(shape.line, self.style["line"]["fill"])
            dict_to_object(self.style["line"], shape.line, exclude=["fill"])

        dict_to_object(self.style["shape_bounds"], shape)
        if "rotation" in dir(shape):
            shape.rotation = self.style["rotation"]
        return shape

    @property
    def img_path(self):
        return self.data[0]

    @img_path.setter
    def img_path(self, img_path: str):
        self.data[0] = img_path

    @property
    def caption(self):
        return self.data[2]

    @caption.setter
    def caption(self, caption: str):
        self.data[2] = caption

    def to_html(self, style_args: StyleArg) -> str:
        if not style_args.show_image:
            return ""
        if not self.caption:
            raise ValueError(
                f"caption not found for picture {self.shape_idx} of slide {self.slide_idx}"
            )
        return (
            self.indent
            + f"<img {self.get_inline_style(style_args)} alt='{self.caption}'/>"
        )


class Placeholder(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: SlidePlaceholder,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        assert (
            sum(
                [
                    shape.has_text_frame,
                    shape.has_chart,
                    shape.has_table,
                    isinstance(shape, PlaceholderPicture),
                ]
            )
            == 1
        ), "placeholder should have only one type"
        if isinstance(shape, PlaceholderPicture):
            data = Picture.from_shape(
                slide_idx,
                shape_idx,
                shape,
                style,
                text_frame,
                config,
                slide_area,
                level,
            )
        elif shape.has_text_frame:
            data = TextBox.from_shape(
                slide_idx,
                shape_idx,
                shape,
                style,
                text_frame,
                config,
                slide_area,
                level,
            )
        else:
            raise ValueError(f"unsupported placeholder {shape.placeholder_type}")
        return data


class GroupShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXGroupShape,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        data = [
            ShapeElement.from_shape(
                slide_idx,
                (shape_idx + 1) * 100 + i,
                sub_shape,
                config,
                slide_area,
                level=level + 1,
            )
            for i, sub_shape in enumerate(shape.shapes)
        ]
        for idx, shape_bounds in enumerate(parse_groupshape(shape)):
            data[idx].style["shape_bounds"] = shape_bounds
        return cls(
            slide_idx, shape_idx, style, data, text_frame, slide_area, level=level
        )

    def build(self, slide: PPTXSlide):
        for shape in self.data:
            shape.build(slide)
        return slide

    def to_pptc(self):
        return "\n".join([shape.to_pptc() for shape in self.data])

    def __iter__(self):
        for shape in self.data:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __eq__(self, __value: object) -> bool:
        if not isinstance(__value, GroupShape) or len(self.data) != len(__value.data):
            return False
        for shape1, shape2 in zip(self.data, __value.data):
            if isinstance(shape1, type(shape2)):
                return False
        return True

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}: {self.data}"

    def to_html(self, style_args: StyleArg) -> str:
        return (
            self.indent
            + f"<div class='{self.group_label}'{self.get_inline_style(style_args)}>\n"
            + "\n".join([shape.to_html(style_args) for shape in self.data])
            + "\n"
            + self.indent
            + "</div>\n"
        )


class FreeShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXAutoShape,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        data = {
            "shape_type": shape.auto_shape_type.real,
            "svg_tag": str(shape.auto_shape_type).split()[0].lower(),
        }
        return cls(
            slide_idx, shape_idx, style, data, text_frame, slide_area, level=level
        )

    def to_html(self, style_args: StyleArg) -> str:
        textframe = self.text_frame.to_html(style_args)
        if not textframe:
            return ""
        return (
            f"{self.indent}<div data-shape-type='{self.data['svg_tag']}'{self.get_inline_style(style_args)}>"
            + f"\n{textframe}"
            + f"\n{self.indent}</div>"
        )


class Connector(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXConnector,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        """
        Convert a connector to a freeform shape.
        """
        return FreeShape(
            slide_idx,
            shape_idx,
            style,
            {"shape_type": "connector", "svg_tag": "connector"},
            text_frame,
            slide_area,
            level,
        )


class SlidePage:
    """
    A class to represent a slide page in a presentation.
    """

    def __init__(
        self,
        shapes: list[ShapeElement],
        slide_idx: int,
        real_idx: int,
        background_xml: str,
        slide_notes: str,
        slide_layout_name: str,
        slide_title: str,
        slide_width: int,
        slide_height: int,
    ):
        self.shapes = shapes
        self.slide_idx = slide_idx
        self.real_idx = real_idx
        self.background_xml = background_xml
        self.slide_notes = slide_notes
        self.slide_layout_name = slide_layout_name
        self.slide_title = slide_title
        self.slide_width = slide_width
        self.slide_height = slide_height
        groups_shapes_labels = []
        for shape in self.shape_filter(GroupShape):
            for group_shape in groups_shapes_labels:
                if group_shape == shape:
                    shape.group_label = group_shape.group_label
                    continue
                groups_shapes_labels.append(shape)
                shape.group_label = f"group_{len(groups_shapes_labels)}"

    @classmethod
    def from_slide(
        cls,
        slide: PPTXSlide,
        slide_idx: int,
        real_idx: int,
        slide_width: int,
        slide_height: int,
        config: Config,
    ):
        """
        Create a SlidePage from a PPTXSlide.

        Args:
            slide (PPTXSlide): The slide object.
            slide_idx (int): The index of the slide.
            real_idx (int): The real index of the slide.
            slide_width (int): The width of the slide.
            slide_height (int): The height of the slide.
            config (Config): The configuration object.

        Returns:
            SlidePage: The created SlidePage.
        """
        shapes = [
            ShapeElement.from_shape(
                slide_idx, i, shape, config, slide_width * slide_height
            )
            for i, shape in enumerate(slide.shapes)
            if shape.visible
        ]
        background_xml = extract_fill(slide.background)
        slide_layout_name = slide.slide_layout.name if slide.slide_layout else None
        slide_title = slide.shapes.title.text if slide.shapes.title else None
        slide_notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )
        return cls(
            shapes,
            slide_idx,
            real_idx,
            background_xml,
            slide_notes,
            slide_layout_name,
            slide_title,
            slide_width,
            slide_height,
        )

    def build(self, slide: PPTXSlide):
        for ph in slide.placeholders:
            ph.element.getparent().remove(ph.element)

        apply_fill(slide.background, self.background_xml)

        for shape in self.shapes:
            build_shape = shape.build(slide)
            for closure in shape.closures:
                try:
                    closure.apply(build_shape)
                except:
                    raise ValueError("Failed to apply closures to slides")
        return slide

    def shape_filter(self, shape_type: type, shapes: list[ShapeElement] = None):
        """
        Filter shapes in the slide by type.

        Args:
            shape_type (type): The type of shapes to filter.
            shapes (list[ShapeElement]): The shapes to filter.

        Yields:
            ShapeElement: The filtered shapes.
        """
        if shapes is None:
            shapes = self.shapes
        for shape in shapes:
            if isinstance(shape, shape_type):
                yield shape
            elif isinstance(shape, GroupShape):
                yield from self.shape_filter(shape_type, shape.data)

    def get_content_type(self):
        """
        Get the content type of the slide.
        """
        if len(list(self.shape_filter(Picture))) > 0:
            return "picture"
        return "text"

    def to_html(self, style_args: StyleArg = None, **kwargs) -> str:
        """
        Represent the slide page in HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.
            **kwargs: Additional arguments.

        Returns:
            str: The HTML representation of the slide page.
        """
        if style_args is None:
            style_args = StyleArg(**kwargs)
        return "".join(
            [
                "<!DOCTYPE html>\n<html>\n",
                (f"<title>{self.slide_title}</title>\n" if self.slide_title else ""),
                f'<body style="width:{self.slide_width}pt; height:{self.slide_height}pt;">\n',
                "\n".join([shape.to_html(style_args) for shape in self.shapes]),
                "</body>\n</html>\n",
            ]
        )

    def to_pptc(self):
        """
        Represent the slide page in PPTC format.
        """
        return "\n".join([shape.to_pptc() for shape in self.shapes])

    def to_text(self, show_image: bool = False) -> str:
        """
        Represent the slide page in text.
        """
        text_content = "\n".join(
            [
                shape.text_frame.text.strip()
                for shape in self.shapes
                if shape.text_frame.is_textframe
            ]
        )
        if show_image:
            for image in self.shape_filter(Picture):
                if not image.caption:
                    raise ValueError(
                        f"caption not found for picture {image.shape_idx} of slide {image.slide_idx}"
                    )
                text_content += "\n" + "Image: " + image.caption
        return text_content

    @property
    def text_length(self):
        """
        Get the length of the text in the slide page.
        """
        return sum([len(shape.text_frame) for shape in self.shapes])

    def __iter__(self):
        for shape in self.shapes:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __len__(self):
        return len(self.shapes)


class Presentation:
    """
    PPTAgent's representation of a presentation.
    Aiming at a more readable and editable interface.
    """

    def __init__(
        self,
        slides: list[SlidePage],
        error_history: list[str],
        slide_width: float,
        slide_height: float,
        file_path: str,
        num_pages: int,
    ) -> None:
        """
        Initialize the Presentation.
        """
        self.slides = slides
        self.error_history = error_history
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.num_pages = num_pages
        self.source_file = file_path
        self.prs = PPTXPre(self.source_file)
        self.layout_mapping = {layout.name: layout for layout in self.prs.slide_layouts}
        self.prs.core_properties.last_modified_by = "PPTAgent"

    @classmethod
    def from_file(cls, file_path: str, config: Config):
        """
        Parse a Presentation from a file.
        """
        prs = PPTXPre(file_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slides = []
        error_history = []
        slide_idx = 0
        layouts = [layout.name for layout in prs.slide_layouts]
        num_pages = len(prs.slides)
        for slide in prs.slides:
            if slide._element.get("show") == "0":
                continue  # will not be printed to pdf

            slide_idx += 1
            try:
                if slide.slide_layout.name not in layouts:
                    raise ValueError(
                        f"slide layout {slide.slide_layout.name} not found"
                    )
                slides.append(
                    SlidePage.from_slide(
                        slide,
                        slide_idx - len(error_history),
                        slide_idx,
                        slide_width.pt,
                        slide_height.pt,
                        config,
                    )
                )
            except Exception as e:
                error_history.append((slide_idx, str(e)))
                if config.DEBUG:
                    print(
                        f"Warning in slide {slide_idx} of {file_path}: {traceback.format_exc()}"
                    )

        return cls(
            slides, error_history, slide_width, slide_height, file_path, num_pages
        )

    def save(self, file_path, layout_only=False):
        """
        Save the presentation to a file.

        Args:
            file_path (str): The file path to save the presentation.
            layout_only (bool): Whether to save only the layout for slide clustering.
        """
        self.clear_slides()
        for slide in self.slides:
            if layout_only:
                self.clear_images(slide.shapes)
            pptx_slide = self.build_slide(slide)
            if layout_only:
                self.clear_text(pptx_slide.shapes)
        self.prs.save(file_path)

    def build_slide(self, slide: SlidePage) -> PPTXSlide:
        """
        Build a slide in the presentation.
        """
        return slide.build(
            self.prs.slides.add_slide(self.layout_mapping[slide.slide_layout_name])
        )

    def clear_slides(self):
        """
        Delete all slides from the presentation.
        """
        while len(self.prs.slides) != 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def clear_images(self, shapes: list[ShapeElement]):
        for shape in shapes:
            if isinstance(shape, GroupShape):
                self.clear_images(shape.data)
            elif isinstance(shape, Picture):
                shape.img_path = "resource/pic_placeholder.png"

    def clear_text(self, shapes: list[BaseShape]):
        for shape in shapes:
            if isinstance(shape, PPTXGroupShape):
                self.clear_text(shape.shapes)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "a" * len(run.text)

    def to_text(self, show_image: bool = False) -> str:
        """
        Represent the presentation in text.
        """
        return "\n----\n".join(
            [
                (
                    f"Slide {slide.slide_idx} of {len(self.prs.slides)}\n"
                    + (f"Title:{slide.slide_title}\n" if slide.slide_title else "")
                    + slide.to_text(show_image)
                )
                for slide in self.slides
            ]
        )

    def __len__(self):
        return len(self.slides)


SHAPECAST: dict[int, ShapeElement] = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: FreeShape,
    MSO_SHAPE_TYPE.PLACEHOLDER: Placeholder,
    MSO_SHAPE_TYPE.PICTURE: Picture,
    MSO_SHAPE_TYPE.GROUP: GroupShape,
    MSO_SHAPE_TYPE.TEXT_BOX: TextBox,
    MSO_SHAPE_TYPE.LINE: Connector,
}

if __name__ == "__main__":
    from copy import deepcopy
    from glob import glob

    config = Config("/tmp")
    presentation = deepcopy(
        Presentation.from_file("runs/pptx/cip_default_template/source.pptx", config)
    ).save("./test.pptx")
    for pptx in glob("data/*/pptx/*/source.pptx"):
        presentation = deepcopy(Presentation.from_file(pptx, config))
        for slide in presentation.slides:
            print(slide.to_html(show_image=False))
            print("\033c", end="")
