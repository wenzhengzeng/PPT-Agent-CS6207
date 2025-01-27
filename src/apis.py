import inspect
import os
import re
import traceback
from copy import deepcopy
from dataclasses import dataclass
from enum import Enum
from functools import partial
from typing import Union

import PIL
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.util import Pt

from presentation import Closure, Picture, ShapeElement, SlidePage
from utils import runs_merge


@dataclass
class HistoryMark:
    """
    Mark the execution status of the API call, comment and a line of code.
    """

    API_CALL_ERROR = "api_call_error"
    API_CALL_CORRECT = "api_call_correct"
    COMMENT_CORRECT = "comment_correct"
    COMMENT_ERROR = "comment_error"
    CODE_RUN_ERROR = "code_run_error"
    CODE_RUN_CORRECT = "code_run_correct"


class CodeExecutor:
    """
    Execute code actions and manage API call history, and providing error feedback.
    """

    def __init__(self, retry_times: int):
        """
        Initialize the CodeExecutor.

        Args:
            retry_times (int): The number of times to retry failed actions.
        """
        self.api_history = []
        self.command_history = []
        self.code_history = []
        self.retry_times = retry_times
        self.registered_functions = API_TYPES.all_funcs()
        self.function_regex = re.compile(r"^[a-z]+_[a-z_]+\(.+\)")

    def get_apis_docs(self, funcs: list[callable], show_example: bool = True) -> str:
        """
        Get the documentation for a list of API functions.

        Args:
            funcs (list[callable]): A list of functions to document.
            show_example (bool): Whether to show examples in the documentation.

        Returns:
            str: The formatted API documentation.
        """
        api_doc = []
        for func in funcs:
            sig = inspect.signature(func)
            params = []
            for name, param in sig.parameters.items():
                if name == "slide":
                    continue
                param_str = name
                if param.annotation != inspect.Parameter.empty:
                    param_str += f": {param.annotation.__name__}"
                if param.default != inspect.Parameter.empty:
                    param_str += f" = {repr(param.default)}"
                params.append(param_str)
            signature = f"def {func.__name__}({', '.join(params)})"
            if not show_example:
                api_doc.append(signature)
                continue
            doc = inspect.getdoc(func)
            if doc is not None:
                signature += f"\n\t{doc}"
            api_doc.append(signature)
        return "\n\n".join(api_doc)

    def execute_actions(
        self, actions: str, edit_slide: SlidePage, found_code: bool = False
    ) -> Union[tuple[str, str], None]:
        """
        Execute a series of actions on a slide.

        Args:
            actions (str): The actions to execute.
            edit_slide (SlidePage): The slide to edit.
            found_code (bool): Whether code was found in the actions.

        Returns:
            tuple: The API lines and traceback if an error occurs.
            None: If no error occurs.
        """
        api_calls = actions.strip().split("\n")
        self.api_history.append(
            [HistoryMark.API_CALL_ERROR, edit_slide.slide_idx, actions]
        )
        for line_idx, line in enumerate(api_calls):
            try:
                if line_idx == len(api_calls) - 1 and not found_code:
                    raise ValueError(
                        "No code block found in the output, please output the api calls without any prefix."
                    )
                if line.startswith("def"):
                    raise PermissionError("The function definition were not allowed.")
                if line.startswith("#"):
                    if len(self.command_history) != 0:
                        self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
                    self.command_history.append([HistoryMark.COMMENT_ERROR, line, None])
                    continue
                if not self.function_regex.match(line):
                    continue
                found_code = True
                func = line.split("(")[0]
                if func not in self.registered_functions:
                    raise NameError(f"The function {func} is not defined.")
                # only one of clone and del can be used in a row
                if func.startswith("clone") or func.startswith("del"):
                    tag = func.split("_")[0]
                    if (
                        self.command_history[-1][-1] == None
                        or self.command_history[-1][-1] == tag
                    ):
                        self.command_history[-1][-1] = tag
                    else:
                        raise ValueError(
                            "Invalid command: Both 'clone_paragraph' and 'del_paragraph'/'del_image' are used within a single command. "
                            "Each command must only perform one of these operations based on the quantity_change."
                        )
                self.code_history.append([HistoryMark.CODE_RUN_ERROR, line, None])
                partial_func = partial(self.registered_functions[func], edit_slide)
                eval(line, {}, {func: partial_func})
                self.code_history[-1][0] = HistoryMark.CODE_RUN_CORRECT
            except:
                trace_msg = traceback.format_exc()
                if len(self.code_history) != 0:
                    self.code_history[-1][-1] = trace_msg
                api_lines = (
                    "\n".join(api_calls[: line_idx - 1])
                    + f"\n--> Error Line: {line}\n"
                    + "\n".join(api_calls[line_idx + 1 :])
                )
                return api_lines, trace_msg
        if len(self.command_history) != 0:
            self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
        self.api_history[-1][0] = HistoryMark.API_CALL_CORRECT


# supporting functions
def element_index(slide: SlidePage, element_id: int) -> ShapeElement:
    """
    Find the an element in a slide.

    Args:
        slide (SlidePage): The slide
        element_id (int): The ID of the element.

    Returns:
        ShapeElement: The shape corresponding to the element ID.

    Raises:
        IndexError: If the element is not found.
    """
    for shape in slide:
        if shape.shape_idx == element_id:
            return shape
    raise IndexError(f"Cannot find element {element_id}, is it deleted or not exist?")


def replace_para(paragraph_id: int, new_text: str, shape: BaseShape):
    """
    Replace the text of a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    runs_merge(para).text = new_text


def clone_para(paragraph_id: int, shape: BaseShape):
    """
    Clone a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    shape.text_frame.paragraphs[-1]._element.addnext(parse_xml(para._element.xml))


def del_para(paragraph_id: int, shape: BaseShape):
    """
    Delete a paragraph from a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    para._element.getparent().remove(para._element)


# api functions
def del_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    """
    Delete a paragraph from a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to delete.

    Raises:
        IndexError: If the paragraph is not found.
    """
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    for para in shape.text_frame.paragraphs:
        if para.idx == paragraph_id:
            shape.text_frame.paragraphs.remove(para)
            shape._closures["delete"].append(
                Closure(partial(del_para, para.real_idx), para.real_idx)
            )
            return
    else:
        raise IndexError(
            f"Cannot find the paragraph {paragraph_id} of the element {div_id},"
            "may refer to a non-existed paragraph or you haven't cloned enough paragraphs beforehand."
        )


def del_image(slide: SlidePage, figure_id: int):
    """
    Delete an image from a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        figure_id (int): The ID of the image to delete.
    """
    shape = element_index(slide, figure_id)
    assert isinstance(shape, Picture), "The element is not a Picture."
    slide.shapes.remove(shape)


def replace_paragraph(slide: SlidePage, div_id: int, paragraph_id: int, text: str):
    """
    Replace the text of a paragraph in a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to replace.
        text (str): The new text to replace with.

    Raises:
        IndexError: If the paragraph is not found.
    """
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    for para in shape.text_frame.paragraphs:
        if para.idx == paragraph_id:
            para.text = text
            shape._closures["replace"].append(
                Closure(
                    partial(replace_para, para.real_idx, text),
                    para.real_idx,
                )
            )
            return
    else:
        raise IndexError(
            f"Cannot find the paragraph {paragraph_id} of the element {div_id},"
            "Please: "
            "1. check if you refer to a non-existed paragraph."
            "2. check if you already deleted it."
        )


def replace_image(slide: SlidePage, img_id: int, image_path: str):
    """
    Replace an image in a slide.

    Args:
        slide (SlidePage): The slide containing the image.
        img_id (int): The ID of the image to replace.
        image_path (str): The path to the new image.

    Raises:
        ValueError: If the image path does not exist.
    """
    if not os.path.exists(image_path):
        raise ValueError(
            f"The image {image_path} does not exist, consider use del_image if image_path in the given command is faked"
        )
    shape = element_index(slide, img_id)
    assert isinstance(shape, Picture), "The element is not a Picture."
    img_size = PIL.Image.open(image_path).size
    r = min(shape.width / img_size[0], shape.height / img_size[1])
    new_width = int(img_size[0] * r)
    new_height = int(img_size[1] * r)
    shape.width = Pt(new_width)
    shape.height = Pt(new_height)
    shape.img_path = image_path


def clone_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    """
    Clone a paragraph in a slide.

    Args:
        slide (SlidePage): The slide containing the paragraph.
        div_id (int): The ID of the division containing the paragraph.
        paragraph_id (int): The ID of the paragraph to clone.

    Raises:
        IndexError: If the paragraph is not found.

    Mention: the cloned paragraph will have a paragraph_id one greater than the current maximum in the parent element.
    """
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    max_idx = max([para.idx for para in shape.text_frame.paragraphs])
    for para in shape.text_frame.paragraphs:
        if para.idx != paragraph_id:
            continue
        shape.text_frame.paragraphs.append(deepcopy(para))
        shape.text_frame.paragraphs[-1].idx = max_idx + 1
        shape.text_frame.paragraphs[-1].real_idx = len(shape.text_frame.paragraphs) - 1
        shape._closures["clone"].append(
            Closure(
                partial(clone_para, para.real_idx),
                para.real_idx,
            )
        )
        return
    raise IndexError(
        f"Cannot find the paragraph {paragraph_id} of the element {div_id}, may refer to a non-existed paragraph."
    )


class API_TYPES(Enum):
    Agent = [
        replace_image,
        del_image,
        clone_paragraph,
        replace_paragraph,
        del_paragraph,
    ]

    @classmethod
    def all_funcs(cls) -> dict[str, callable]:
        funcs = {}
        for attr in dir(cls):
            if attr.startswith("__"):
                continue
            funcs |= {func.__name__: func for func in getattr(cls, attr).value}
        return funcs


if __name__ == "__main__":
    print(CodeExecutor(0).get_apis_docs(API_TYPES.Agent.value))
